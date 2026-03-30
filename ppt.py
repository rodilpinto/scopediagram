from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt


SLIDE_WIDTH = Emu(12192000)
SLIDE_HEIGHT = Emu(6858000)


class Palette:
    background = RGBColor(255, 255, 255)
    title = RGBColor(41, 52, 43)
    muted = RGBColor(84, 98, 90)
    line = RGBColor(196, 203, 197)
    header_green = RGBColor(50, 106, 62)
    header_orange = RGBColor(222, 136, 37)
    objective_fill = RGBColor(252, 241, 221)
    regulators_fill = RGBColor(244, 234, 210)
    resources_fill = RGBColor(235, 242, 229)
    event_fill = RGBColor(243, 245, 241)
    inputs_fill = RGBColor(241, 247, 237)
    activities_fill = RGBColor(255, 245, 230)
    outputs_fill = RGBColor(244, 250, 240)
    lane_label_fill = RGBColor(85, 128, 70)
    lane_label_text = RGBColor(255, 255, 255)


def _set_slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = Palette.background


def _set_text_style(paragraph, size: int, *, bold: bool = False, color: RGBColor | None = None, align=PP_ALIGN.LEFT) -> None:
    paragraph.alignment = align
    for run in paragraph.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color


def _add_textbox(slide, left: int, top: int, width: int, height: int, text: str, *, size: int, bold: bool = False, color: RGBColor | None = None, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    _set_text_style(paragraph, size=size, bold=bold, color=color, align=align)
    return shape


def _add_panel(slide, left: int, top: int, width: int, height: int, fill: RGBColor):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = Palette.line
    shape.line.width = Pt(1)
    shape.text_frame.margin_left = Pt(10)
    shape.text_frame.margin_right = Pt(10)
    shape.text_frame.margin_top = Pt(8)
    shape.text_frame.margin_bottom = Pt(6)
    return shape


def _set_panel_text(shape, title: str, body: str, *, body_size: int = 12) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.vertical_anchor = MSO_ANCHOR.TOP

    p1 = frame.paragraphs[0]
    p1.text = title
    _set_text_style(p1, size=12, bold=True, color=Palette.muted)

    p2 = frame.add_paragraph()
    p2.text = body or "N/A"
    _set_text_style(p2, size=body_size, color=Palette.title)


def _load_logo_path() -> str | None:
    candidates = [
        Path("assets/logo-camara.png"),
        Path("assets/logo_camara.png"),
        Path("assets/camara_logo.png"),
        Path("assets/camara-dos-deputados.png"),
    ]
    for path in candidates:
        if path.exists():
            return str(path)
    return None


def _add_header(slide, title: str) -> None:
    header = _add_panel(slide, 220000, 140000, 11750000, 620000, RGBColor(249, 251, 248))
    header.line.color.rgb = Palette.line
    logo_path = _load_logo_path()
    if logo_path:
        slide.shapes.add_picture(logo_path, Emu(320000), Emu(220000), height=Emu(320000))
        text_left = 1700000
    else:
        emblem = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Emu(320000), Emu(250000), Emu(260000), Emu(260000))
        emblem.fill.solid()
        emblem.fill.fore_color.rgb = Palette.header_green
        emblem.line.color.rgb = Palette.header_green
        orange = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Emu(600000), Emu(250000), Emu(140000), Emu(260000))
        orange.fill.solid()
        orange.fill.fore_color.rgb = Palette.header_orange
        orange.line.color.rgb = Palette.header_orange
        text_left = 820000

    _add_textbox(slide, text_left, 220000, 2800000, 200000, "Câmara dos Deputados", size=15, bold=True, color=Palette.header_green)
    _add_textbox(slide, text_left, 400000, 2800000, 160000, "Secretaria de Controle Interno", size=11, color=Palette.muted)
    _add_textbox(slide, 7200000, 260000, 4400000, 220000, title, size=20, bold=True, color=Palette.title, align=PP_ALIGN.RIGHT)


def _vertical_lane(slide, left: int, top: int, width: int, height: int, label: str, items: list[str], fill: RGBColor) -> None:
    panel = _add_panel(slide, left, top, width, height, fill)
    panel.text = ""

    label_length = height - 320000
    label_height = 420000
    label_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Emu(left - (label_length // 2) + 160000),
        Emu(top + (height // 2) - (label_height // 2)),
        Emu(label_length),
        Emu(label_height),
    )
    label_box.fill.solid()
    label_box.fill.fore_color.rgb = Palette.lane_label_fill
    label_box.line.color.rgb = Palette.lane_label_fill

    label_frame = label_box.text_frame
    label_frame.clear()
    label_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    label_frame.word_wrap = True
    label_paragraph = label_frame.paragraphs[0]
    label_paragraph.text = label
    _set_text_style(label_paragraph, size=11, bold=True, color=Palette.lane_label_text, align=PP_ALIGN.CENTER)
    label_box.rotation = 270

    content_box = slide.shapes.add_textbox(Emu(left + 380000), Emu(top + 120000), Emu(width - 500000), Emu(height - 240000))
    content_frame = content_box.text_frame
    content_frame.clear()
    content_frame.word_wrap = True
    content_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    content_frame.margin_left = 0
    content_frame.margin_right = 0
    content_frame.margin_top = 0
    content_frame.margin_bottom = 0
    for index, item in enumerate(items or ["N/A"]):
        paragraph = content_frame.paragraphs[0] if index == 0 else content_frame.add_paragraph()
        paragraph.text = item
        paragraph.bullet = True
        _set_text_style(paragraph, size=11, color=Palette.title)


def _lines(items: list[str]) -> str:
    return "\n".join(items) if items else "N/A"


def _build_scope_slide(prs: Presentation, title: str, objective: str, regulators: list[str], left_items: list[str], middle_label: str, middle_items: list[str], right_items: list[str], resources: list[str], start_event: str, end_event: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide)
    _add_header(slide, title)

    objective_shape = _add_panel(slide, 280000, 960000, 11600000, 760000, Palette.objective_fill)
    _set_panel_text(objective_shape, "OBJETIVO", objective, body_size=12)

    regulators_shape = _add_panel(slide, 280000, 1880000, 11600000, 1160000, Palette.regulators_fill)
    _set_panel_text(regulators_shape, "REGULADORES", _lines(regulators), body_size=11)

    _vertical_lane(slide, 280000, 3260000, 3520000, 2600000, "ENTRADAS", left_items, Palette.inputs_fill)
    _vertical_lane(slide, 3960000, 3260000, 3520000, 2600000, middle_label, middle_items, Palette.activities_fill)
    _vertical_lane(slide, 7640000, 3260000, 4200000, 2600000, "SAÍDAS", right_items, Palette.outputs_fill)

    resources_shape = _add_panel(slide, 280000, 6060000, 11600000, 760000, Palette.resources_fill)
    _set_panel_text(resources_shape, "RECURSOS DE SUPORTE", _lines(resources), body_size=11)

    start_shape = _add_panel(slide, 280000, 6600000, 5600000, 520000, Palette.event_fill)
    _set_panel_text(start_shape, "EVENTO DE INÍCIO", start_event, body_size=11)
    end_shape = _add_panel(slide, 6280000, 6600000, 5600000, 520000, Palette.event_fill)
    _set_panel_text(end_shape, "EVENTO DE FIM", end_event, body_size=11)


def _build_cover_slide(prs: Presentation, scope) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide)
    _add_header(slide, scope.process.name)
    hero = _add_panel(slide, 1600000, 1900000, 9000000, 2100000, Palette.objective_fill)
    _set_panel_text(hero, "DIAGRAMAS DE ESCOPO", scope.process.objective, body_size=18)
    _add_textbox(slide, 2200000, 4300000, 7800000, 600000, scope.process.name.upper(), size=24, bold=True, color=Palette.title, align=PP_ALIGN.CENTER)


def generate_ppt(scope):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _build_cover_slide(prs, scope)
    _build_scope_slide(
        prs,
        f"Processo {scope.process.name}",
        scope.process.objective,
        scope.global_elements.regulators if scope.global_elements else [],
        scope.global_elements.inputs if scope.global_elements else [],
        "SUBPROCESSOS",
        [subprocess.name for subprocess in scope.subprocesses],
        scope.global_elements.outputs if scope.global_elements else [],
        scope.global_elements.resources if scope.global_elements else [],
        scope.process.start_event,
        scope.process.end_event,
    )

    for subprocess in scope.subprocesses:
        _build_scope_slide(
            prs,
            subprocess.name,
            subprocess.objective or scope.process.objective,
            subprocess.regulators or (scope.global_elements.regulators if scope.global_elements else []),
            subprocess.inputs,
            "ATIVIDADES",
            subprocess.activities,
            subprocess.outputs,
            subprocess.resources or (scope.global_elements.resources if scope.global_elements else []),
            subprocess.start_event or scope.process.start_event,
            subprocess.end_event or scope.process.end_event,
        )

    return prs


def generate_ppt_bytes(scope) -> bytes:
    prs = generate_ppt(scope)
    output = BytesIO()
    prs.save(output)
    return output.getvalue()
